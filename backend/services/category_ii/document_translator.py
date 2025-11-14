# backend/services/document_translator.py
import os
import sys
import re
import pdfplumber
import fitz  # PyMuPDF - pentru editare PDF păstrând formatarea
from deep_translator import GoogleTranslator
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import simpleSplit
from PIL import Image
import pytesseract
import requests 
from services.progress_bar import send_pages_progress
GEMINI_API_KEY = "AIzaSyCMr_9mKOj2lJkq-sfA47cTMmCRwO3Jj0U"
GEMINI_MODEL = "gemini-2.0-flash-lite" 
# BASE = "http://86.126.134.77:11434/api/generate"
# La început în document_translator.py, adaugă:

# Fix encoding pentru Windows
if sys.platform == 'win32':
    import codecs
    sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer)
    sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer)
# Path către font Arial (vine cu Windows)
FONT_PATH = r'C:\Windows\Fonts\arial.ttf'
# Configurare Tesseract pentru Windows
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

PROCESSED_FOLDER = "processed"
os.makedirs(PROCESSED_FOLDER, exist_ok=True)

def translate_text_gemini(text, src_lang='ru', dest_lang='ro'):
    """
    Traduce text folosind Google Gemini API
    """
    if not text or not text.strip():
        return text
    
    # Mapping limbi
    lang_names = {
        'ru': 'Russian',
        'en': 'English', 
        'zh': 'Chinese',
        'ja': 'Japanese',
        'ro': 'Romanian'
    }
    
    src_name = lang_names.get(src_lang, src_lang)
    dest_name = lang_names.get(dest_lang, dest_lang)
    
    # Prompt pentru Gemini
    prompt = f"""Translate this {src_name} text to {dest_name}. 

Rules:
- Provide ONLY the translation, nothing else
- Do NOT add explanations, notes, or comments
- Preserve formatting (bullet points, numbers, line breaks)
- Use proper Romanian grammar and diacritics

Text to translate:
{text}

Translation:"""
    
    # URL API Gemini
    api_url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent?key={GEMINI_API_KEY}"
    
    # Payload
    payload = {
        "contents": [{
            "parts": [{
                "text": prompt
            }]
        }],
        "generationConfig": {
            "temperature": 0.1,
            "topK": 1,
            "topP": 0.8,
            "maxOutputTokens": 2048,
        }
    }
    
    try:
        print(f"[Gemini] Traducere {src_lang} → {dest_lang}: '{text[:50]}...'")
        
        response = requests.post(
            api_url,
            json=payload,
            headers={'Content-Type': 'application/json'},
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            translated = result['candidates'][0]['content']['parts'][0]['text'].strip()
            
            # Curăță răspunsul
            import re
            prefixes = [r'^Translation:\s*', r'^Romanian translation:\s*', r'^Traducere:\s*']
            for pattern in prefixes:
                translated = re.sub(pattern, '', translated, flags=re.IGNORECASE)
            
            translated = re.sub(r'\n{3,}', '\n\n', translated)
            translated = re.sub(r' {2,}', ' ', translated).strip()
            
            print(f"[Gemini] ✅ Rezultat: '{translated[:50]}...'")
            return translated
        else:
            error_msg = response.json().get('error', {}).get('message', f'HTTP {response.status_code}')
            print(f"[Gemini] ❌ Eroare API: {error_msg}")
            raise Exception(f"Gemini API error: {error_msg}")
            
    except Exception as e:
        print(f"[Gemini] ❌ Eroare: {e}")
        raise

def _clean_Gemini_contamination(text, original_text):
    """
    🧹 Elimină contaminarea cu prompt-uri și instrucțiuni din răspunsul Gemeni
    """
    import re
    
    # Liste de cuvinte cheie care indică contaminare
    contamination_patterns = [
        r'^Translation:\s*',  # "Translation: ..."
        r'^Romanian translation:\s*',  # "Romanian translation: ..."
        r'^Traducere:\s*',  # "Traducere: ..."
        r'^Reguli[e]? de respectat:.*?(?=\n[A-Z]|\n•|\n\d|\Z)',  # Blochează "Regulile de respectat"
        r'RULES TO FOLLOW:.*?(?=\n[A-Z]|\n•|\n\d|\Z)',  # Blochează "RULES TO FOLLOW"
        r'^-\s*Nu oferi.*?(?=\n[A-Z]|\n•|\n\d|\Z)',  # Blochează instrucțiuni
        r'^-\s*Provide ONLY.*?(?=\n[A-Z]|\n•|\n\d|\Z)',
        r'^-\s*Folosește.*?(?=\n[A-Z]|\n•|\n\d|\Z)',
        r'^-\s*Conjugări.*?(?=\n[A-Z]|\n•|\n\d|\Z)',
        r'^-\s*Diacritice.*?(?=\n[A-Z]|\n•|\n\d|\Z)',
    ]
    
    cleaned = text
    
    # Elimină pattern-urile de contaminare
    for pattern in contamination_patterns:
        cleaned = re.sub(pattern, '', cleaned, flags=re.MULTILINE | re.DOTALL | re.IGNORECASE)
    
    # Elimină linii goale multiple
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    
    # 🔍 VERIFICARE SUPLIMENTARĂ: Dacă textul conține "Regulile" sau "RULES", elimină tot până la primul paragraf real
    if 'eguli' in cleaned or 'RULE' in cleaned.upper():
        # Găsește primul paragraf după instrucțiuni (începe cu majusculă sau bullet/număr)
        match = re.search(r'(?:^|\n)([A-ZĂÂÎȘȚ•\d].*)', cleaned, re.MULTILINE)
        if match:
            cleaned = match.group(1)
    
    # Elimină spații multiple
    cleaned = re.sub(r' {2,}', ' ', cleaned)
    
    cleaned = cleaned.strip()
    
    # 🚨 VERIFICARE FINALĂ: Dacă traducerea e ciudată (prea scurtă/prea lungă), raportează
    original_words = len(original_text.split())
    translated_words = len(cleaned.split())
    
    ratio = translated_words / max(original_words, 1)
    
    if ratio > 3:  # Traducerea e de 3x mai lungă = suspect
        print(f"⚠️ WARNING: Traducere suspectă (ratio: {ratio:.1f}x)")
        print(f"   Original: {original_words} cuvinte")
        print(f"   Tradus: {translated_words} cuvinte")
        print(f"   Output: '{cleaned[:100]}...'")
    
    return cleaned

def test_ollama_cleaning():
    """Test pentru curățarea contaminării Ollama"""
    
    test_cases = [
        # Caz 1: Contaminare cu "Regulile de respectat"
        (
            "Original text",
            """Regulile de respectat:
 - Nu oferi nicio altă explicație
 - Folosește gramatica perfectă
 Traducere: Regatele mari ale oamenilor erau...""",
            "Regatele mari ale oamenilor erau..."
        ),
        
        # Caz 2: Prefix "Translation:"
        (
            "Original",
            "Translation: Textul tradus corect",
            "Textul tradus corect"
        ),
        
        # Caz 3: Traducere curată (nu trebuie modificată)
        (
            "Original", 
            "• Orașele din Elderath au căzut",
            "• Orașele din Elderath au căzut"
        )
    ]
    
    print("\n🧪 === TESTING Gemini CLEANING ===")
    for i, (original, contaminated, expected) in enumerate(test_cases, 1):
        result = _clean_Gemini_contamination(contaminated, original)
        status = "✅" if expected in result else "❌"
        print(f"\nTest {i}: {status}")
        print(f"  Input: '{contaminated[:60]}...'")
        print(f"  Output: '{result[:60]}...'")
        print(f"  Expected: '{expected[:60]}...'")
    
    print("\n=== END TESTING ===\n")

def translate_text(text, src_lang='en', dest_lang='ro'):
    """Traduce text cu mapping explicit pentru cuvinte cunoscute"""
    if not text or not text.strip():
        return text
    # MAPPING EXPLICIT pentru titluri și cuvinte cunoscute
    known_translations = {
        'ru': {
            'Привет': 'Salut',
            'Здравствуйте': 'Bună ziua',
            'Спасибо': 'Mulțumesc',
            'Пожалуйста': 'Cu plăcere'
        },
        'en': {
            'Hello': 'Salut',
            'Good morning': 'Bună dimineața',
            'Thank you': 'Mulțumesc'
        },
        'ja': {
            'こんにちは': 'Salut',
            'ありがとう': 'Mulțumesc'
        },
        'zh': {
            '你好': 'Salut',
            '谢谢': 'Mulțumesc'
        }
    }
    # Verifică dacă avem traducere cunoscută
    if src_lang in known_translations:
        if text.strip() in known_translations[src_lang]:
            translated = known_translations[src_lang][text.strip()]
            print(f"[Translate] ✅ Mapping direct: '{text}' -> '{translated}'")
            return translated
    try:
        # 🆕 ÎNCEARCĂ OLLAMA PRIMUL
        return translate_text_gemini(text, src_lang, dest_lang)
        
    except Exception as gemini_error:
        # Fallback la Google Translate
        print(f"[Translate] ⚠️ Gemini failed, using Google Translate: {gemini_error}")
        
        try:
            translator = GoogleTranslator(source=src_lang, target=dest_lang)
            
            if len(text) > 4500:
                chunks = [text[i:i+4500] for i in range(0, len(text), 4500)]
                translated = ' '.join([translator.translate(chunk) for chunk in chunks])
            else:
                translated = translator.translate(text)
            
            translated = _fix_word_spacing(translated)
            translated = _fix_romanian_translation_errors(translated)
            return translated
            
        except Exception as e:
            print(f"[Translate] ❌ EROARE: {e}")
            raise
def _fix_word_spacing(text):
    """Adaugă spații între cuvinte lipite după traducere"""
    import re
    # Caută pattern-uri de tipul: MajusculăMajusculă -> Majusculă Majusculă
    # Ex: "TreileaIdee" -> "Treia Idee"
    text = re.sub(r'([a-zăâîșț])([A-ZĂÂÎȘȚ])', r'\1 \2', text)
    return text
def _fix_romanian_translation_errors(text):
    """
    Repară erorile comune ale Google Translate pentru română
    """
    if not text:
        return text
    
    # 1. Repară spațiile greșite în jurul cratimelor
    # "s -a" -> "s-a", "într -un" -> "într-un", etc.
    text = re.sub(r'\b(\w+)\s+-\s*(\w+)\b', r'\1-\2', text)
    
    # 2. Repară punct dublu
    text = re.sub(r'\.\.+', '.', text)
    
    # 3. Repară acordul greșit
    common_errors = {
        'ardeau. ': 'ardea. ',
        'erau alunecoase': 'erau alunecoase',
        'pietruirea erau': 'pietrele erau',
        'strălucea Crimson': 'strălucea roșu-aprins',
    }
    
    for wrong, correct in common_errors.items():
        text = text.replace(wrong, correct)
    
    # 4. Repară spațiile excesive
    text = re.sub(r' {2,}', ' ', text)
    
    return text.strip()
def _detect_numbered_list(text):
    """
    Detectează dacă textul începe cu o numerotare de listă.
    Suportă: 1. 2. 3. | 1) 2) 3) | a. b. c. | i. ii. iii.
    """
    patterns = [
        r'^\s*\d+\.\s',          # 1. 2. 3. (cu spații opționale)
        r'^\s*\d+\)\s',          # 1) 2) 3)
        r'^\s*[a-zA-Z]\.\s',     # a. b. c.
        r'^\s*[a-zA-Z]\)\s',     # a) b) c)
        r'^\s*[ivxIVX]+\.\s',    # i. ii. iii.
        r'^\s*[ivxIVX]+\)\s',    # i) ii) iii)
        r'^\s*\(\d+\)\s',        # (1) (2) (3)
        r'^\s*\([a-zA-Z]\)\s',   # (a) (b) (c)
    ]
    
    stripped = text.strip()
    for pattern in patterns:
        if re.match(pattern, stripped):
            return True
    return False

def _detect_title_improved(text):
    """Detectare ÎMBUNĂTĂȚITĂ pentru titluri"""
    stripped = text.strip()
    # Lista titlurilor EXACTE cunoscute
    exact_titles = [
        'Привет', 'Hello', 'Salut', 'Bună ziua',
        'Introduction', 'Introducere', 'Введение',
        'Chapter', 'Capitolul', 'Глава'
    ]
    if stripped in exact_titles:
        return True
    # Criterii generale
    is_short = len(stripped) < 50
    is_single_line = '\n' not in stripped
    max_words = len(stripped.split()) <= 4
    no_end_punct = not stripped.endswith(('.', ',', ';', ':', '!', '?'))
    not_bullet = not (stripped.startswith('•') or stripped.startswith('â€¢'))
    # Cuvinte comune care NU apar în titluri
    common_words_ru = ['и', 'или', 'но', 'для', 'от', 'с', 'в', 'на', 'из']
    common_words_ro = ['în', 'și', 'sau', 'dar', 'pentru', 'de', 'la', 'cu', 'a']
    common_words_en = ['and', 'or', 'but', 'for', 'from', 'with', 'in', 'on', 'at', 'the', 'a', 'an']
    all_common = common_words_ru + common_words_ro + common_words_en
    has_common_words = any(word in stripped.lower().split() for word in all_common)
    return (
        is_short and 
        is_single_line and 
        max_words and 
        no_end_punct and 
        not_bullet and 
        not has_common_words
    )
def translate_image_text(image_path, src_lang='en', dest_lang='ro'):
    """Extrage text din imagine și îl traduce"""
    try:
        img = Image.open(image_path)
        # OCR cu pytesseract
        text = pytesseract.image_to_string(img, lang='eng')
        if text.strip():
            translated = translate_text(text, src_lang, dest_lang)
            return translated
        return ""
    except Exception as e:
        print(f"[Image OCR Error] {e}")
        return ""
# ==================== WORD (DOCX) ====================
def translate_word_document(input_path, output_path=None, src_lang='en', dest_lang='ro'):
    """Traduce document Word păstrând formatarea"""
    if output_path is None:
        base = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(PROCESSED_FOLDER, f"{base}_RO.docx")
    print(f"[Word] Traducere document: {input_path}")
    doc = Document(input_path)
    translated_paragraphs = 0
    # Traduce paragrafe
    for para in doc.paragraphs:
        if para.text.strip():
            original_text = para.text
            translated_text = translate_text(original_text, src_lang, dest_lang)
            # Păstrează formatarea - înlocuiește textul în primul run
            if para.runs:
                para.runs[0].text = translated_text
                # Șterge celelalte run-uri
                for run in para.runs[1:]:
                    run.text = ''
            else:
                para.text = translated_text
            
            translated_paragraphs += 1
            if translated_paragraphs % 10 == 0:
                print(f"  Traduse {translated_paragraphs} paragrafe...")
    
    # Traduce tabele
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    cell.text = translate_text(cell.text, src_lang, dest_lang)
    
    # Salvează documentul tradus
    doc.save(output_path)
    print(f"[Word] ✅ Salvat: {output_path}")
    
    return {
        'output_path': output_path,
        'translated_paragraphs': translated_paragraphs,
        'status': 'success'
    }

# ==================== PowerPoint (PPTX) ====================
def translate_ppt_document(input_path, output_path=None, src_lang='en', dest_lang='ro'):
    """Traduce prezentare PowerPoint păstrând formatarea"""
    if output_path is None:
        base = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(PROCESSED_FOLDER, f"{base}_RO.pptx")
    
    print(f"[PPT] Traducere prezentare: {input_path}")
    
    prs = Presentation(input_path)
    translated_count = 0
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"  Slide {slide_num}/{len(prs.slides)}")
        
        for shape in slide.shapes:
            # Traduce text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            original = run.text
                            run.text = translate_text(original, src_lang, dest_lang)
                            translated_count += 1
            
            # Traduce tabele
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            cell.text = translate_text(cell.text, src_lang, dest_lang)
                            translated_count += 1
    
    prs.save(output_path)
    print(f"[PPT] ✅ Salvat: {output_path}")
    
    return {
        'output_path': output_path,
        'translated_items': translated_count,
        'total_slides': len(prs.slides),
        'status': 'success'
    }

def debug_pdf_extraction(input_path, max_pages=1):
    """
    Debug pentru a vedea EXACT ce text extrage pdfplumber din PDF
    """
    print("\n" + "="*60)
    print("🔍 DEBUG: EXTRAGERE TEXT DIN PDF")
    print("="*60)
    
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages[:max_pages]):
            print(f"\n📄 Pagina {page_num + 1}:")
            print("-" * 60)
            
            # Extrage text cu layout
            text = page.extract_text(layout=True, x_tolerance=3, y_tolerance=3)
            
            # Afișează primele 500 caractere RAW
            print("RAW TEXT (primele 500 char):")
            print(repr(text[:500]))
            print()
            
            # Verifică probleme comune
            print("⚠️ PROBLEME DETECTATE:")
            
            # 1. Encoding corupt
            if 'Ð' in text or 'Ñ' in text:
                print("  ❌ Encoding CORUPT detectat (Ð, Ñ)")
            
            # 2. Cuvinte lipite (2+ cuvinte fără spațiu)
            words = text.split()
            long_words = [w for w in words if len(w) > 30]
            if long_words:
                print(f"  ⚠️ Cuvinte FOARTE LUNGI (posibil lipite): {len(long_words)}")
                for w in long_words[:3]:
                    print(f"    - {w[:50]}...")
            
            # 3. Cuvinte rupte cu hyphen
            if '-\n' in text:
                print("  ⚠️ Cuvinte RUPTE pe linii (hyphen)")
            
            # 4. Spații multiple consecutive
            if '  ' in text:
                print("  ⚠️ Spații MULTIPLE consecutive")
            
            print()
            
            # Afișează primele 3 paragrafe pentru verificare
            paragraphs = [p for p in text.split('\n\n') if p.strip()][:3]
            print("📝 PRIMELE 3 PARAGRAFE:")
            for i, para in enumerate(paragraphs, 1):
                print(f"\n  Paragraf {i}:")
                print(f"  {para[:200]}...")
    
    print("\n" + "="*60)
# ==================== PDF ====================
def clean_extracted_text(text):
    """
    Curăță textul extras din PDF FĂRĂ să strice cratimele din cuvinte
    """
    if not text:
        return text
    
    # 1. Repară DOAR cuvinte rupte cu hyphen la sfârșit de linie
    #    Exemplu: "захват-\nчиков" -> "захватчиков"
    #    DAR NU atinge "s-a" sau "și-a" care sunt corecte!
    text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
    
    # 2. Înlocuiește spații multiple cu unul singur
    #    DAR păstrează spațiile în "s -a" pentru că le reparăm după
    text = re.sub(r' {2,}', ' ', text)
    
    # 3. Înlocuiește multiple newlines cu maximum 2
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 4. Elimină spații la început/sfârșit de linii
    lines = []
    for line in text.split('\n'):
        lines.append(line.strip())
    text = '\n'.join(lines)
    
    # 5. Elimină caractere de control invizibile (dar păstrează \n, \t)
    text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F-\x9F]', '', text)
    
    # 6. REPARĂ spațiile greșite în jurul cratimelor
    #    "s -a" -> "s-a", "și -a" -> "și-a", etc.
    text = re.sub(r'\b(\w+)\s+-\s*(\w+)\b', r'\1-\2', text)
    
    return text
# ==================== PDF SIMPLIFICAT ====================
def translate_pdf_document(input_path, output_path=None, src_lang='en', dest_lang='ro'):
    """
    Traduce PDF RECREÂND documentul de la ZERO (SOLUȚIE FINALĂ)
    """
    
    if output_path is None:
        base = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(PROCESSED_FOLDER, f"{base}_RO.pdf")
    
    print(f"\n{'='*60}")
    print(f"[PDF] Recreează PDF CU TRADUCERE (de la ZERO)")
    print(f"[PDF] {input_path} → {output_path}")
    print(f"{'='*60}\n")
    
    # Deschide PDF ORIGINAL pentru extragere
    doc_original = fitz.open(input_path)
    
    # Creează PDF NOU (gol)
    doc_new = fitz.open()
    
    total_blocks = 0
    translated_blocks = 0
    
    # Procesează fiecare pagină
    for page_num in range(len(doc_original)):
        
       
        page_original = doc_original[page_num]
        print(f"📄 Pagina {page_num + 1}/{len(doc_original)}")
        
        # Creează PAGINĂ NOUĂ în PDF-ul nou (aceeași dimensiune)
        page_new = doc_new.new_page(
            width=page_original.rect.width,
            height=page_original.rect.height
        )
        
        # COPIAZĂ IMAGINILE din original (dacă există)
        try:
            image_list = page_original.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc_original.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Găsește poziția imaginii
                for item in page_original.get_image_bbox(xref):
                    rect = fitz.Rect(item)
                    # Inserează imaginea în pagina nouă
                    page_new.insert_image(rect, stream=image_bytes)
                    print(f"   🖼️  Copiat imagine #{img_index}")
                    break
        except Exception as e:
            print(f"   ⚠️ Nu pot copia imagini: {e}")
        
        # Extrage BLOCURI TEXT din original
        blocks = page_original.get_text("dict")["blocks"]
        
        # Sortează după Y
        text_blocks = []
        for block in blocks:
            if block.get("type") == 0:
                bbox = block.get("bbox")
                if bbox:
                    text_blocks.append(block)
        
        text_blocks.sort(key=lambda b: b["bbox"][1])
        print(f"   📦 {len(text_blocks)} blocuri text")
        
        # Procesează fiecare bloc
        for block_idx, block in enumerate(text_blocks):
            total_blocks += 1
            
            # Extrage text
            block_text = ""
            font_size = 11
            
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    block_text += span.get("text", "") + " "
                    if span.get("size"):
                        font_size = span["size"]
            
            block_text = block_text.strip()
            
            if len(block_text) < 2:
                continue
            
            bbox = block["bbox"]
            rect = fitz.Rect(bbox)
            
            try:
                # Traduce
                translated_text = translate_text(block_text, src_lang, dest_lang)
                
                if not translated_text:
                    translated_text = block_text  # Fallback la original
                
                # SCRIE în PAGINA NOUĂ (pe PDF curat)
                result = page_new.insert_textbox(
                    rect,
                    translated_text,
                    fontname="helv",
                    fontsize=font_size,
                    color=(0, 0, 0),
                    align=fitz.TEXT_ALIGN_LEFT
                )
                
                if result > 0:
                    translated_blocks += 1
                    print(f"   ✅ #{block_idx}: '{translated_text[:30]}'")
                else:
                    # Reduce font
                    result = page_new.insert_textbox(
                        rect,
                        translated_text,
                        fontname="helv",
                        fontsize=font_size * 0.85,
                        color=(0, 0, 0),
                        align=fitz.TEXT_ALIGN_LEFT
                    )
                    
                    if result > 0:
                        translated_blocks += 1
                        print(f"   ✅ #{block_idx}: '{translated_text[:30]}' (font redus)")
                    else:
                        # Mai mic
                        result = page_new.insert_textbox(
                            rect,
                            translated_text,
                            fontname="helv",
                            fontsize=font_size * 0.7,
                            color=(0, 0, 0),
                            align=fitz.TEXT_ALIGN_LEFT
                        )
                        
                        if result > 0:
                            translated_blocks += 1
                            print(f"   ✅ #{block_idx}: '{translated_text[:30]}' (font mic)")
                        else:
                            print(f"   ❌ #{block_idx}: Nu încape textul")
            
            except Exception as e:
                print(f"   ❌ Eroare: {e}")
                continue
    send_pages_progress(done=page_num + 1, total=len(doc_original))
    
    # Închide originalul
    doc_original.close()
    
    # Salvează PDF-ul NOU
    total_pages = len(doc_new)
    
    try:
        doc_new.save(output_path, garbage=4, deflate=True, clean=True)
        doc_new.close()
    except Exception as e:
        print(f"[PDF] ❌ Eroare salvare: {e}")
        doc_new.close()
        raise
    
    print(f"\n{'='*60}")
    print(f"[PDF] ✅ SUCCES!")
    print(f"[PDF] PDF NOU creat cu {total_pages} pagini")
    print(f"[PDF] Traduse: {translated_blocks}/{total_blocks} blocuri")
    print(f"[PDF] Salvat: {output_path}")
    print(f"{'='*60}\n")
    
    return {
        'output_path': output_path,
        'total_pages': total_pages,
        'translated_blocks': translated_blocks,
        'status': 'success'
    }


def _parse_pdf_paragraphs_improved(text):
    """
    Parsare CORECTĂ care păstrează bullet points ȘI liste numerotate EXACT ca în original.
    """
    lines = text.split('\n')
    structure = []
    current_paragraph = []
    
    print(f"    📄 Analizez {len(lines)} linii pentru structură EXACTĂ...")
    
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        
        if not stripped:  # Linie goală
            if current_paragraph:
                # Salvează paragraful curent
                para_text = '\n'.join(current_paragraph).strip()
                if para_text:
                    is_bullet = para_text.startswith('•') or para_text.startswith('•')
                    is_numbered = _detect_numbered_list(para_text)
                    is_title = _detect_title_improved(para_text) and not is_bullet and not is_numbered
                    
                    structure.append({
                        'type': 'paragraph',
                        'text': para_text,
                        'is_title': is_title,
                        'is_bullet': is_bullet,
                        'is_numbered': is_numbered,
                        'indent_level': 0,
                        'line_count': len(current_paragraph)
                    })
                    
                    type_label = '📌 TITLU' if is_title else ('• BULLET' if is_bullet else ('🔢 NUMBERED' if is_numbered else 'NORMAL'))
                    print(f"      {type_label}: '{para_text[:50]}...'")
                
                current_paragraph = []
            
            # Adaugă spațiu gol
            structure.append({'type': 'empty_line'})
            i += 1
        
        # =====================================================
        # VERIFICARE LISTE NUMEROTATE - ÎNAINTE DE BULLETS!
        # =====================================================
        elif _detect_numbered_list(stripped):
            # NUMBERED LIST DETECTION
            if current_paragraph:
                # Finalizează paragraful anterior
                para_text = '\n'.join(current_paragraph).strip()
                if para_text:
                    structure.append({
                        'type': 'paragraph',
                        'text': para_text,
                        'is_title': False,
                        'is_bullet': False,
                        'is_numbered': False,
                        'indent_level': 0,
                        'line_count': len(current_paragraph)
                    })
                current_paragraph = []
                structure.append({'type': 'empty_line'})
            
            # Colectează ÎNTREG item-ul numerotat
            numbered_lines = [line]
            i += 1
            
            # Continuă să colecteze până la:
            # 1. Următorul item numerotat SAU
            # 2. Bullet point SAU  
            # 3. Linie goală urmată de text non-listă
            while i < len(lines):
                next_line = lines[i]
                next_stripped = next_line.strip()
                
                # Stop la următorul item numerotat sau bullet
                if _detect_numbered_list(next_stripped):
                    break
                if next_stripped.startswith('•') or next_stripped.startswith('•'):
                    break
                
                # Stop la linie goală DOAR dacă urmează paragraf nou
                if not next_stripped:
                    if i + 1 < len(lines):
                        line_after_empty = lines[i + 1].strip()
                        # Verifică dacă e continuare sau paragraf nou
                        if not line_after_empty or (line_after_empty and not _detect_numbered_list(line_after_empty) and not line_after_empty.startswith('•')):
                            # Verifică indentare pentru continuare
                            if i + 1 < len(lines) and len(lines[i + 1]) > len(lines[i + 1].lstrip()) + 2:
                                pass  # E indentat, continuă
                            else:
                                break  # Paragraf nou, stop
                
                numbered_lines.append(next_line)
                i += 1
            
            # Creează item numerotat complet
            numbered_text = '\n'.join(numbered_lines).strip()
            structure.append({
                'type': 'paragraph',
                'text': numbered_text,
                'is_title': False,
                'is_bullet': False,
                'is_numbered': True,
                'indent_level': 0,
                'line_count': len(numbered_lines)
            })
            
            print(f"      🔢 NUMBERED COMPLET: '{numbered_text[:60]}...' ({len(numbered_lines)} linii)")
        
        # =====================================================
        # BULLET POINTS - LOGICA EXISTENTĂ RĂMÂNE NESCHIMBATĂ
        # =====================================================
        elif stripped.startswith('•') or stripped.startswith('•'):
            # BULLET POINT DETECTION - colectează TOT conținutul
            if current_paragraph:
                # Finalizează paragraful anterior
                para_text = '\n'.join(current_paragraph).strip()
                if para_text:
                    structure.append({
                        'type': 'paragraph',
                        'text': para_text,
                        'is_title': False,
                        'is_bullet': False,
                        'is_numbered': False,
                        'indent_level': 0,
                        'line_count': len(current_paragraph)
                    })
                current_paragraph = []
                structure.append({'type': 'empty_line'})
            
            # Colectează ÎNTREG bullet point-ul
            bullet_lines = [line]
            i += 1
            
            # Continuă să colecteze linii PÂNĂ LA:
            # 1. Următorul bullet SAU
            # 2. Linie goală urmată de text non-bullet SAU
            # 3. Două linii goale consecutive
            while i < len(lines):
                next_line = lines[i]
                next_stripped = next_line.strip()
                
                # Stop la următorul bullet
                if next_stripped.startswith('•') or next_stripped.startswith('•'):
                    break
                
                # Stop la următoarea listă numerotată
                if _detect_numbered_list(next_stripped):
                    break
                
                # Stop la linie goală DOAR dacă următoarea linie e:
                # - goală și ea (sfârșit paragraf) SAU
                # - text care NU începe cu spații (paragraf nou)
                if not next_stripped:
                    if i + 1 < len(lines):
                        line_after_empty = lines[i + 1].strip()
                        # Dacă următoarea e și ea goală SAU e text nou (nu continuare), STOP
                        if not line_after_empty or (line_after_empty and not line_after_empty.startswith('•')):
                            # Verifică dacă e continuare indentată
                            if i + 1 < len(lines) and len(lines[i + 1]) > len(lines[i + 1].lstrip()) + 2:
                                # E indentată, continuă
                                pass
                            else:
                                break
                
                bullet_lines.append(next_line)
                i += 1
            
            # Creează bullet point complet
            bullet_text = '\n'.join(bullet_lines).strip()
            structure.append({
                'type': 'paragraph',
                'text': bullet_text,
                'is_title': False,
                'is_bullet': True,
                'is_numbered': False,
                'indent_level': 0,
                'line_count': len(bullet_lines)
            })
            
            print(f"      • BULLET COMPLET: '{bullet_text[:60]}...' ({len(bullet_lines)} linii)")
        
        # =====================================================
        # TEXT NORMAL - COLECTEAZĂ ÎN PARAGRAF CURENT
        # =====================================================
        else:
            # Linie normală
            current_paragraph.append(line)
            i += 1
    
    # =====================================================
    # ULTIMUL PARAGRAF
    # =====================================================
    if current_paragraph:
        para_text = '\n'.join(current_paragraph).strip()
        if para_text:
            is_bullet = para_text.startswith('•') or para_text.startswith('•')
            is_numbered = _detect_numbered_list(para_text)
            is_title = _detect_title_improved(para_text) and not is_bullet and not is_numbered
            
            structure.append({
                'type': 'paragraph',
                'text': para_text,
                'is_title': is_title,
                'is_bullet': is_bullet,
                'is_numbered': is_numbered,
                'indent_level': 0,
                'line_count': len(current_paragraph)
            })
    
    # =====================================================
    # STATISTICI FINALE
    # =====================================================
    print(f"    📊 Structură finală: {len(structure)} elemente")
    
    # Debug statistics
    bullets = [s for s in structure if s.get('is_bullet')]
    numbered = [s for s in structure if s.get('is_numbered')]
    titles = [s for s in structure if s.get('is_title')]
    
    print(f"    🔵 Bullet points: {len(bullets)}")
    print(f"    🔢 Liste numerotate: {len(numbered)}")
    print(f"    📌 Titluri: {len(titles)}")
    
    return structure

# AJUSTEAZĂ funcția de creare PDF pentru spațiere compactă:

# ÎNLOCUIEȘTE și funcția de creare PDF cu spațiere corectă:

#def _is_title_line(text):
    #"""Detectează dacă textul este un titlu"""
   # stripped = text.strip()
    
    # Titlu dacă:
    # - Este scurt (< 60 caractere)
    # - Este pe o singură linie
    # - Nu conține punct la sfârșit
    # - Este centrat sau are caractere speciale
    
    return (
        len(stripped) < 60 and 
        '\n' not in stripped and
        not stripped.endswith('.') and
        (len(stripped.split()) <= 3 or stripped.isupper())
    )

def _get_indent_level(line):
    """Calculează nivelul de indentare al unei linii"""
    if not line:
        return 0
    return len(line) - len(line.lstrip())

def _create_pdf_with_improved_structure(all_page_structures, output_path):
    """Crează PDF cu formatare IDENTICĂ pentru bullets ȘI liste numerotate"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os
    
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    
    # Font setup
    font_regular = 'Helvetica'
    font_bold = 'Helvetica-Bold'
    
    try:
        if os.path.exists(r'C:\Windows\Fonts\arial.ttf'):
            pdfmetrics.registerFont(TTFont('Arial', r'C:\Windows\Fonts\arial.ttf'))
            font_regular = 'Arial'
            if os.path.exists(r'C:\Windows\Fonts\arialbd.ttf'):
                pdfmetrics.registerFont(TTFont('Arial-Bold', r'C:\Windows\Fonts\arialbd.ttf'))
                font_bold = 'Arial-Bold'
            print(f"[PDF] ✅ Font Arial încărcat")
    except Exception as e:
        print(f"[PDF] ⚠️ Folosim Helvetica: {e}")
    
    for page_num, page_structure in enumerate(all_page_structures):
        if not page_structure:
            c.showPage()
            continue
        
        y_position = height - 40
        
        for elem_idx, element in enumerate(page_structure):
            if element['type'] == 'paragraph':
                text = element['text']
                is_title = element.get('is_title', False)
                is_bullet = element.get('is_bullet', False)
                is_numbered = element.get('is_numbered', False)  # Tratare pentru liste numerotate
                
                if not text.strip():
                    continue
                
                # Înlătură line breaks pentru re-wrapping corect
                text = ' '.join(text.split())
                
                # Configurare stil
                if is_title:
                    font_size = 14
                    font_name = font_bold
                    line_spacing = 17
                    center_text = True
                    left_margin = 60
                    right_margin = 60
                    use_indent = False
                elif is_bullet:
                    font_size = 11
                    font_name = font_regular
                    line_spacing = 14
                    center_text = False
                    left_margin = 90
                    right_margin = 60
                    use_indent = False
                    
                    if not text.startswith('•'):
                        text = '• ' + text
                elif is_numbered:  # Stil pentru liste numerotate - IDENTIC cu bullets
                    font_size = 11
                    font_name = font_regular
                    line_spacing = 14
                    center_text = False
                    left_margin = 90  # Același indent ca bullets
                    right_margin = 60
                    use_indent = False
                    
                    # NU adăugăm nimic - păstrăm numerotarea originală
                else:
                    # Paragraf normal cu alineat
                    font_size = 11
                    font_name = font_regular
                    line_spacing = 14
                    center_text = False
                    left_margin = 60
                    right_margin = 60
                    use_indent = True
                
                # Setează fontul pentru calcul lățime
                c.setFont(font_name, font_size)
                
                # Calculează lățimea maximă disponibilă
                if use_indent:
                    max_width_first = width - (left_margin + 30) - right_margin
                    max_width_rest = width - left_margin - right_margin
                else:
                    max_width_first = width - left_margin - right_margin
                    max_width_rest = max_width_first
                
                # Wrap text
                lines = _wrap_text_by_width(text, c, font_name, font_size, max_width_rest)
                
                # Dacă prima linie cu alineat e prea lungă, re-împarte
                if use_indent and lines:
                    first_line_width = c.stringWidth(lines[0], font_name, font_size)
                    if first_line_width > max_width_first:
                        lines = _wrap_text_by_width(text, c, font_name, font_size, max_width_first)
                
                # Desenează liniile
                for line_num, line in enumerate(lines):
                    if y_position < 60:
                        c.showPage()
                        y_position = height - 40
                    
                    try:
                        c.setFont(font_name, font_size)
                        
                        if center_text and is_title:
                            text_width = c.stringWidth(line, font_name, font_size)
                            x_position = (width - text_width) / 2
                        else:
                            # Alineat pentru prima linie paragraf normal
                            if use_indent and line_num == 0:
                                x_position = left_margin + 30
                            # Indent pentru continuare bullets și numbered
                            elif (is_bullet or is_numbered) and line_num > 0:
                                x_position = left_margin + 15
                            else:
                                x_position = left_margin
                        
                        c.drawString(x_position, y_position, line)
                        
                    except Exception as e:
                        print(f"      ⚠️ Eroare: {e}")
                        c.setFont('Helvetica', font_size)
                        c.drawString(left_margin, y_position, line)
                    
                    y_position -= line_spacing
                
                # Spațiere între elemente
                next_is_bullet_or_numbered = False
                if elem_idx + 1 < len(page_structure):
                    next_elem = page_structure[elem_idx + 1]
                    if next_elem['type'] == 'paragraph':
                        next_is_bullet_or_numbered = next_elem.get('is_bullet') or next_elem.get('is_numbered')
                
                if is_title:
                    y_position -= 18
                elif is_bullet or is_numbered:
                    y_position -= 10
                else:
                    if next_is_bullet_or_numbered:
                        y_position -= 8
                    else:
                        y_position -= 10
            
            elif element['type'] == 'empty_line':
                prev_elem = page_structure[elem_idx - 1] if elem_idx > 0 else None
                next_elem = page_structure[elem_idx + 1] if elem_idx + 1 < len(page_structure) else None
                
                # Spațiere inteligentă între elemente
                if (prev_elem and prev_elem['type'] == 'paragraph' and 
                    not prev_elem.get('is_bullet') and not prev_elem.get('is_numbered') and
                    next_elem and next_elem['type'] == 'paragraph' and 
                    (next_elem.get('is_bullet') or next_elem.get('is_numbered'))):
                    y_position -= 6
                else:
                    y_position -= 10
        
        c.showPage()
    
    c.save()
    print(f"[PDF] ✅ PDF generat cu structură COMPLETĂ - bullets ȘI liste numerotate!")


def _wrap_text_by_width(text, canvas_obj, font_name, font_size, max_width):
    """Împarte textul bazat pe lățime REALĂ în PDF, nu pe număr de caractere"""
    if not text or not text.strip():
        return []
    
    words = text.split()
    lines = []
    current_line = []
    
    for word in words:
        # Testează linia cu noul cuvânt
        test_line = ' '.join(current_line + [word])
        test_width = canvas_obj.stringWidth(test_line, font_name, font_size)
        
        # Dacă depășește lățimea ȘI avem deja cuvinte în linie, începe linie nouă
        if test_width > max_width and current_line:
            lines.append(' '.join(current_line))
            current_line = [word]
        else:
            current_line.append(word)
    
    # Adaugă ultima linie
    if current_line:
        lines.append(' '.join(current_line))
    
    return lines
    

# ÎNLOCUIEȘTE și funcția de detectare titlu pentru mai multă precizie

def _wrap_text_smart(text, max_chars):
    """Împarte textul INTELIGENT fără să taie cuvinte la mijloc"""
    if len(text) <= max_chars:
        return [text]
    
    words = text.split()
    lines = []
    current_line = []
    current_length = 0
    
    for word in words:
        # Calculează lungimea cu spațiu
        word_with_space = len(word) + (1 if current_line else 0)
        test_length = current_length + word_with_space
        
        # CRITICAL: Nu adăuga cuvântul dacă depășește limita ȘI linia are deja cuvinte
        if test_length > max_chars and current_line:
            lines.append(' '.join(current_line))
            current_line = [word]
            current_length = len(word)
        else:
            current_line.append(word)
            current_length = test_length
    
    # Adaugă ultima linie
    if current_line:
        lines.append(' '.join(current_line))
    
    return lines

# ADAUGĂ și o funcție pentru debugging în document_translator.py
def debug_pdf_structure(input_path):
    """Funcție pentru a vedea cum se parsează structura PDF-ului"""
    print("🔍 DEBUG: Analizez structura PDF...")
    
    with pdfplumber.open(input_path) as pdf:
        for page_num, page in enumerate(pdf.pages[:1]):  # doar prima pagină
            text = page.extract_text(layout=True)
            print(f"\n📄 Pagina {page_num + 1}:")
            print("=" * 50)
            print(repr(text[:500]))  # afișează caracterele speciale
            print("=" * 50)
            
            structure = _parse_pdf_paragraphs(text)
            print(f"\n🏗️ Structura detectată ({len(structure)} elemente):")
            for i, elem in enumerate(structure[:10]):  # primele 10 elemente
                if elem['type'] == 'paragraph':
                    title_flag = "📌 TITLU" if elem.get('is_title') else ""
                    bullet_flag = "• BULLET" if elem.get('is_bullet') else ""
                    print(f"  #{i}: {title_flag} {bullet_flag} | '{elem['text'][:40]}...'")
                else:
                    print(f"  #{i}: SPAȚIU GOL") 

# ADAUGĂ și funcția lipsă pentru fallback
def _create_pdf_simple_clean(pages_text, output_path):
    """Fallback simplu pentru când structura detaliată eșuează"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    
    for page_text in pages_text:
        if not page_text or not page_text.strip():
            c.showPage()
            continue
            
        y_position = height - 55
        
        # Împarte în paragrafe
        paragraphs = page_text.split('\n\n')
        
        for para in paragraphs:
            if not para.strip():
                y_position -= 20
                continue
            
            # Wrap text la 90 caractere
            lines = []
            for line in para.split('\n'):
                while len(line) > 90:
                    lines.append(line[:90])
                    line = line[90:]
                if line:
                    lines.append(line)
            
            # Desenează liniile
            for line in lines:
                if y_position < 60:
                    c.showPage()
                    y_position = height - 55
                
                try:
                    c.setFont('Helvetica', 10)
                    c.drawString(60, y_position, line[:100])
                except:
                    pass
                
                y_position -= 15
            
            # Spațiu între paragrafe
            y_position -= 10
        
        c.showPage()
    
    c.save()
    print(f"[PDF] ✅ Fallback clean folosit")

def _parse_page_structure(text):
    """
    Analizează structura textului (paragrafe, titluri, indentări).
    Returnează listă de elemente structurate.
    """
    lines = text.split('\n')
    structure = []
    current_paragraph = []
    
    for line in lines:
        stripped = line.strip()
        
        # Linie goală = separator de paragrafe
        if not stripped:
            if current_paragraph:
                # Salvează paragraful curent
                para_text = ' '.join(current_paragraph)
                
                # Detectează titlu (text scurt, centrat sau cu caractere speciale)
                is_title = (
                    len(para_text) < 50 and 
                    (para_text.startswith('-') or para_text.isupper() or para_text.endswith('-'))
                )
                
                structure.append({
                    'type': 'paragraph',
                    'text': para_text,
                    'indent': 0,
                    'is_title': is_title
                })
                current_paragraph = []
            
            # Adaugă linie goală
            structure.append({'type': 'empty_line'})
        else:
            # Adaugă la paragraful curent
            current_paragraph.append(stripped)
    
    # Ultimul paragraf
    if current_paragraph:
        para_text = ' '.join(current_paragraph)
        is_title = (
            len(para_text) < 50 and 
            (para_text.startswith('-') or para_text.isupper() or para_text.endswith('-'))
        )
        structure.append({
            'type': 'paragraph',
            'text': para_text,
            'indent': 0,
            'is_title': is_title
        })
    
    return structure

def _create_pdf_with_structure(pages_structure, output_path):
    """Creează PDF păstrând structura (paragrafe, titluri, spații)"""
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    
    # Path către font DejaVu
    font_path = os.path.join(os.path.dirname(__file__), '..', 'fonts', 'ttf', 'DejaVuSans.ttf')
    font_path = os.path.abspath(font_path)
    
    font_name = 'Times-Roman'  # default fallback
    
    # Încearcă să încarci DejaVu
    try:
        if os.path.exists(font_path):
            pdfmetrics.registerFont(TTFont('DejaVu', font_path))
            font_name = 'DejaVu'
            print(f"[PDF] ✅ Font DejaVu încărcat")
        else:
            print(f"[PDF] ⚠️  Folosim Times-Roman (diacritice parțiale)")
    except Exception as e:
        print(f"[PDF] ⚠️  Eroare font: {e}")
    
    # Creează document
    doc = SimpleDocTemplate(
        output_path, 
        pagesize=A4,
        rightMargin=2*cm, 
        leftMargin=2*cm,
        topMargin=0.5*cm, 
        bottomMargin=2*cm
    )
    
    # Stiluri pentru diferite tipuri de text
    styles = getSampleStyleSheet()
    
    # Stil pentru paragrafe normale
    style_normal = ParagraphStyle(
        'Normal',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=11,
        leading=16,
        alignment=TA_JUSTIFY,
        spaceAfter=6,
        firstLineIndent=0.5*cm,  # Indent prima linie
    )
    
    # Stil pentru titluri
    style_title = ParagraphStyle(
        'Title',
        parent=styles['Normal'],
        fontName=font_name,
        fontSize=12,
        leading=18,
        alignment=TA_CENTER,
        spaceAfter=12,
        spaceBefore=0,
    )
    
    # Construiește conținut
    story = []
    
    for page_structure in pages_structure:
        if not page_structure:
            story.append(PageBreak())
            continue
        
        for item in page_structure:
            if item['type'] == 'paragraph':
                text = item['text']
                
                # Escape caractere XML
                text = text.replace('&', '&amp;')
                text = text.replace('<', '&lt;')
                text = text.replace('>', '&gt;')
                
                try:
                    # Alege stilul în funcție de tip
                    if item.get('is_title'):
                        para = Paragraph(text, style_title)
                        
                    else:
                        para = Paragraph(text, style_normal)
                    
                    story.append(para)
                    
                except Exception as e:
                    print(f"    ⚠️  Skip: {str(e)[:30]}")
                    continue
            
            elif item['type'] == 'empty_line':
                # Adaugă spațiu între paragrafe
                story.append(Spacer(1, 0.5*cm))
        
        # Pagină nouă
        story.append(PageBreak())
    
    # Generează PDF
    try:
        doc.build(story)
        print(f"[PDF] ✅ PDF generat cu structură păstrată")
    except Exception as e:
        print(f"[PDF] ⚠️  Eroare: {e}")
        # Fallback simplu
        _create_pdf_simple_structure(pages_structure, output_path)

def _create_pdf_simple_structure(pages_structure, output_path):
    """Fallback simplu cu păstrare paragraf"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    
    # Încearcă font DejaVu
    font_path = os.path.join(os.path.dirname(__file__), '..', 'fonts', 'ttf', 'DejaVuSans.ttf')
    try:
        if os.path.exists(font_path):
            pdfmetrics.registerFont(TTFont('DejaVu', font_path))
            font_name = 'DejaVu'
        else:
            font_name = 'Helvetica'
    except:
        font_name = 'Helvetica'
    
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    
    for page_structure in pages_structure:
        y_position = height - 50
        
        for item in page_structure:
            if item['type'] == 'paragraph':
                text = item['text']
                is_title = item.get('is_title', False)
                
                # Wrap text
                max_width = 90 if not is_title else 70
                lines = []
                words = text.split()
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    if len(test_line) <= max_width:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Desenează liniile
                for line in lines:
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    
                    try:
                        if is_title:
                            c.setFont(font_name, 12)
                            # Centrează titlul
                            text_width = c.stringWidth(line, font_name, 12)
                            x_position = (width - text_width) / 2
                            c.drawString(x_position, y_position, line)
                        else:
                            c.setFont(font_name, 10)
                            c.drawString(60, y_position, line)
                    except:
                        pass
                    
                    y_position -= 15
                
                # Spațiu după paragraf
                y_position -= 10
            
            elif item['type'] == 'empty_line':
                y_position -= 15
        
        c.showPage()
    
    c.save()
    print(f"[PDF] ✅ Fallback cu structură folosit")
def _create_pdf_ascii_safe(pages_text, output_path):
    """
    Fallback: Convertește caracterele speciale românești 
    la versiunea lor ASCII-safe pentru fonturi limitate
    """
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    
    def romanize(text):
        """Convertește diacritice românești la echivalente ASCII"""
        replacements = {
            'ă': 'a', 'Ă': 'A',
            'â': 'a', 'Â': 'A', 
            'î': 'i', 'Î': 'I',
            'ș': 's', 'Ș': 'S',
            'ț': 't', 'Ț': 'T'
        }
        for ro, en in replacements.items():
            text = text.replace(ro, en)
        return text
    
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    
    for page_text in pages_text:
        y_position = height - 50
        
        # Convertește textul la ASCII-safe
        page_text = romanize(page_text)
        
        # Split în linii
        lines = []
        for paragraph in page_text.split('\n\n'):
            para_lines = paragraph.split('\n')
            for line in para_lines:
                # Wrap text la 90 caractere
                while len(line) > 90:
                    lines.append(line[:90])
                    line = line[90:]
                if line:
                    lines.append(line)
            lines.append('')  # Spațiu între paragrafe
        
        # Desenează liniile
        for line in lines[:60]:  # Max 60 linii per pagină
            if y_position < 50:
                c.showPage()
                y_position = height - 50
            
            try:
                c.drawString(50, y_position, line[:100])
            except:
                pass
            
            y_position -= 15
        
        c.showPage()
    
    c.save()
    print(f"[PDF] ✅ Fallback ASCII-safe folosit (fără diacritice)")

def _create_pdf_simple(pages_text, output_path):
    """Fallback: PDF simplu dacă reportlab dă erori"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    
    for page_text in pages_text:
        y_position = height - 50
        
        # Split în linii
        lines = []
        for paragraph in page_text.split('\n\n'):
            # Wrap text la 90 caractere
            para_lines = paragraph.split('\n')
            for line in para_lines:
                while len(line) > 90:
                    lines.append(line[:90])
                    line = line[90:]
                if line:
                    lines.append(line)
            lines.append('')  # Spațiu între paragrafe
        
        # Desenează liniile
        for line in lines[:60]:  # Max 60 linii per pagină
            if y_position < 50:
                c.showPage()
                y_position = height - 50
            
            try:
                c.drawString(50, y_position, line[:100])
            except:
                pass  # Skip caractere problematice
            
            y_position -= 15
        
        c.showPage()
    
    c.save()
    print(f"[PDF] ✅ Fallback simplu folosit")

# ==================== FUNCȚIE PRINCIPALĂ ====================
def translate_document(filepath, src_lang='en', dest_lang='ro'):
    """
    Funcție principală pentru traducerea documentelor.
    Suportă: .docx, .pdf, .pptx
    """
    ext = filepath.rsplit('.', 1)[-1].lower()
    
    if ext == 'docx':
        return translate_word_document(filepath, src_lang=src_lang, dest_lang=dest_lang)
    elif ext == 'pptx':
        return translate_ppt_document(filepath, src_lang=src_lang, dest_lang=dest_lang)
    elif ext == 'pdf':
        return translate_pdf_document(filepath, src_lang=src_lang, dest_lang=dest_lang)
    else:
        raise ValueError(f"Format nesuportat: {ext}")
    # ADAUGĂ ACESTE FUNCȚII LA SFÂRȘITUL document_translator.py
# ATENȚIE: Fără spații în plus la început - indentare la nivel 0!

def _clean_text_spacing(text):
    """Curata spatierea excesiva din text"""
    import re
    
    # Inlocuieste spatiile multiple cu unul singur
    text = re.sub(r' +', ' ', text)
    
    # Inlocuieste newline-urile multiple cu maximum 2
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    
    # Elimina spatiile de la inceput si sfarsit de linii
    lines = []
    for line in text.split('\n'):
        lines.append(line.strip())
    
    return '\n'.join(lines)

def _smart_text_wrap(text, max_chars):
    """Imparte textul in linii, respectand cuvintele"""
    words = text.split()
    lines = []
    current_line = []
    current_length = 0
    
    for word in words:
        # Verifica daca adaugarea cuvantului depaseste limita
        if current_length + len(word) + 1 > max_chars and current_line:
            lines.append(' '.join(current_line))
            current_line = [word]
            current_length = len(word)
        else:
            current_line.append(word)
            current_length += len(word) + (1 if current_line else 0)
    
    if current_line:
        lines.append(' '.join(current_line))
    
    return lines

def _convert_romanian_to_ascii(text):
    """Converteste diacriticele romanesti la echivalente ASCII"""
    if not text:
        return text
    
    replacements = {
        'ă': 'a', 'Ă': 'A',
        'â': 'a', 'Â': 'A',
        'î': 'i', 'Î': 'I', 
        'ș': 's', 'Ș': 'S',
        'ț': 't', 'Ț': 'T'
    }
    
    result = text
    for romanian, ascii_char in replacements.items():
        result = result.replace(romanian, ascii_char)
    
    return result

def _clean_text_advanced(text):
        """Curățare avansată a textului pentru spațiere perfectă"""
        import re
    
        # 1. Elimină spațiile multiple
        text = re.sub(r' +', ' ', text)
    
        # 2. Elimină spațiile de la început și sfârșit de linii
        lines = []
        for line in text.split('\n'):
            cleaned = line.strip()
            if cleaned:  # doar liniile cu conținut
                lines.append(cleaned)
    
        # 3. Reconstruiește cu o singură linie goală între paragrafe
        result = '\n'.join(lines)
    
        # 4. Elimină caracterele problematice
        result = re.sub(r'[^\w\s\-•.,!?():;/\\""«»„"—–]', '', result, flags=re.UNICODE)
    
        return result.strip()


def _split_text_preserving_structure(text, max_chars):
        """Împarte textul păstrând structura și bullet points"""
        if len(text) <= max_chars:
            return [text]
        
        # Dacă e bullet point, păstrează • pe prima linie
        if text.strip().startswith('•'):
            bullet_text = text.strip()[1:].strip()  # elimină •
            lines = _wrap_text_smart(bullet_text, max_chars - 2)  # -2 pentru "• "
            if lines:
                lines[0] = '• ' + lines[0]  # adaugă • la prima linie
                # Indentează liniile următoare
                for i in range(1, len(lines)):
                    lines[i] = '  ' + lines[i]  # 2 spații pentru aliniere
            return lines
        else:
            return _wrap_text_smart(text, max_chars)

def _clean_text_advanced(text):
    """Curățare avansată a textului pentru spațiere perfectă"""
    import re
    
    # 1. Elimină spațiile multiple
    text = re.sub(r' +', ' ', text)
    
    # 2. Elimină spațiile de la început și sfârșit de linii
    lines = []
    for line in text.split('\n'):
        cleaned = line.strip()
        if cleaned:  # doar liniile cu conținut
            lines.append(cleaned)
    
    # 3. Reconstruiește cu o singură linie goală între paragrafe
    result = '\n'.join(lines)
    
    # 4. Elimină caracterele problematice
    result = re.sub(r'[^\w\s\-•.,!?():;/\\""«»„"—–]', '', result, flags=re.UNICODE)
    
    return result.strip()

def _split_text_preserving_structure(text, max_chars):
    """Împarte textul păstrând structura și bullet points"""
    if len(text) <= max_chars:
        return [text]
    
    # Dacă e bullet point, păstrează • pe prima linie
    if text.strip().startswith('•'):
        bullet_text = text.strip()[1:].strip()  # elimină •
        lines = _wrap_text_smart(bullet_text, max_chars - 2)  # -2 pentru "• "
        if lines:
            lines[0] = '• ' + lines[0]  # adaugă • la prima linie
            # Indentează liniile următoare
            for i in range(1, len(lines)):
                lines[i] = '  ' + lines[i]  # 2 spații pentru aliniere
        return lines
    else:
        return _wrap_text_smart(text, max_chars)
    
def _split_multiple_bullets(text):
    """
    Împarte un text care conține bullet points multiple în paragrafe separate
    """
    if '•' not in text:
        return [text]
    
    # Dacă textul începe cu bullet, procesează direct
    if text.strip().startswith('•'):
        # Caută toate pozițiile unde apar bullet points
        parts = []
        current_pos = 0
        
        while True:
            # Găsește următorul bullet point
            next_bullet = text.find('•', current_pos + 1)
            
            if next_bullet == -1:  # Nu mai sunt bullet points
                # Adaugă restul textului
                remaining = text[current_pos:].strip()
                if remaining:
                    parts.append(remaining)
                break
            else:
                # Extrage partea curentă
                current_part = text[current_pos:next_bullet].strip()
                if current_part:
                    parts.append(current_part)
                current_pos = next_bullet
        
        return parts if parts else [text]
    
    else:
        # Textul nu începe cu bullet, dar conține bullet points
        # Split la primul bullet
        first_bullet = text.find('•')
        if first_bullet == -1:
            return [text]
        
        before_bullet = text[:first_bullet].strip()
        bullet_part = text[first_bullet:].strip()
        
        result = []
        if before_bullet:
            result.append(before_bullet)
        
        # Procesează partea cu bullet points
        bullet_parts = _split_multiple_bullets(bullet_part)
        result.extend(bullet_parts)
        
        return result

# TESTARE: Adaugă această funcție pentru debugging
def test_bullet_splitting():
    """Test pentru separarea bullet points"""
    test_cases = [
        "• First bullet • Second bullet",
        "Normal text • First bullet • Second bullet",
        "• Single bullet only",
        "No bullets here"
    ]
    
    print("\n=== TESTING BULLET SPLITTING ===")
    for i, test in enumerate(test_cases):
        print(f"Test {i+1}: '{test}'")
        result = _split_multiple_bullets(test)
        for j, part in enumerate(result):
            print(f"  Part {j+1}: '{part}'")
        print()

# Pentru debugging, adaugă temporar în translate_pdf_document:
# test_bullet_splitting()  # Uncomment pentru a testa

# VERIFICARE SUPLIMENTARĂ: Funcție pentru a vedea exact ce extrage pdfplumber
def debug_extraction_method(input_path):
    """Vezi cum pdfplumber extrage textul"""
    import pdfplumber
    
    print("\n=== DEBUG EXTRACTION ===")
    
    with pdfplumber.open(input_path) as pdf:
        page = pdf.pages[0]
        
        # Metoda 1: layout=True
        text1 = page.extract_text(layout=True)
        print("Layout=True:")
        print(repr(text1[:500]))
        print()
        
        # Metoda 2: layout=False  
        text2 = page.extract_text(layout=False)
        print("Layout=False:")
        print(repr(text2[:500]))
        print()
        
        # Metoda 3: extract_words pentru pozițiile exacte
        words = page.extract_words()
        bullet_words = [w for w in words if '•' in w['text']]
        print(f"Bullet words găsite: {len(bullet_words)}")
        for w in bullet_words[:5]:
            print(f"  '{w['text']}' la poziția ({w['x0']}, {w['top']})")

def debug_bullet_parsing(text):
    """Debug pentru a vedea cum se parsează bullet points"""
    lines = text.split('\n')
    print("\n=== DEBUG BULLET PARSING ===")
    
    in_bullet = False
    bullet_content = []
    
    for i, line in enumerate(lines):
        stripped = line.strip()
        if stripped.startswith('•'):
            if in_bullet:
                print(f"Bullet anterior: '{' '.join(bullet_content)[:60]}...'")
            print(f"LINIA {i}: BULLET START -> '{stripped[:60]}...'")
            in_bullet = True
            bullet_content = [stripped]
        elif in_bullet and stripped:
            print(f"LINIA {i}: BULLET CONTINUĂ -> '{stripped[:60]}...'")
            bullet_content.append(stripped)
        elif in_bullet and not stripped:
            print(f"LINIA {i}: BULLET STOP (linie goală)")
            print(f"Bullet complet: '{' '.join(bullet_content)[:80]}...'")
            in_bullet = False
            bullet_content = []
        else:
            print(f"LINIA {i}: normal -> '{stripped[:60]}...'")
    
    if in_bullet:
        print(f"Bullet final: '{' '.join(bullet_content)[:80]}...'")
    
    print("=== END DEBUG ===\n")


