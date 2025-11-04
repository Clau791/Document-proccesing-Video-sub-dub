
from pathlib import Path
import json, zipfile
from datetime import datetime

try:
    from pptx import Presentation  # optional: python-pptx
except Exception:
    Presentation = None

class PPTAnalyzer:
    """
    Minimal PowerPoint analyzer.
    - If python-pptx is available, extracts slide count and text length.
    - Otherwise, falls back to counting slide XML parts in .pptx zip.
    Produces a JSON summary saved into the 'processed' folder.
    """
    def __init__(self, processed_dir: str = "processed"):
        self.processed_dir = Path(processed_dir)
        self.processed_dir.mkdir(parents=True, exist_ok=True)

    def analyze(self, filepath: str):
        path = Path(filepath)
        summary = {
            "filename": path.name,
            "slides": None,
            "text_chars": 0,
            "generated_at": datetime.now().isoformat(),
            "notes": []
        }

        if Presentation is not None:
            try:
                prs = Presentation(str(path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text or "")
                summary["slides"] = len(prs.slides)
                joined = "\n".join(texts)
                summary["text_chars"] = len(joined)
            except Exception as e:
                summary["notes"].append(f"pptx parse fallback: {e}")

        if summary["slides"] is None and path.suffix.lower() == ".pptx":
            try:
                with zipfile.ZipFile(str(path), "r") as z:
                    slides = [n for n in z.namelist() if n.startswith("ppt/slides/slide")]
                    summary["slides"] = len(slides)
                    summary["notes"].append("Counted slides via zip fallback.")
            except Exception as e:
                summary["notes"].append(f"zip fallback failed: {e}")

        if summary["slides"] is None:
            summary["slides"] = 0

        outname = f"{path.stem}_analysis.json"
        outpath = self.processed_dir / outname
        outpath.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")

        return {
            "output_file": f"/download/{outname}",
            "summary": summary
        }
