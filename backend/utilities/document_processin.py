from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import os
import uuid
from datetime import datetime
import json, sys, io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

# Setează encoding-ul pentru Windows
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

app = Flask(_name_)
CORS(app)  # Enable CORS for all routes

# Configuration
UPLOAD_FOLDER = 'uploads'
PROCESSED_FOLDER = 'processed'
IMAGE_FOLDER = os.path.join(PROCESSED_FOLDER, 'images') # Folder pentru imaginile extrase

ALLOWED_EXTENSIONS = {
    'document-parse': {'pdf', 'doc', 'docx', 'ppt', 'pptx', 'jpg', 'jpeg', 'png'},
    'translation': {'pdf', 'doc', 'docx', 'txt'},
    'audio-translation': {'mp3', 'wav', 'm4a', 'ogg'},
    'video-translation': {'mp4', 'avi', 'mov', 'mkv'},
    'video-subtitle': {'mp4', 'avi', 'mov', 'mkv'},
    'video-audio-replace': {'mp4', 'avi', 'mov', 'mkv'},
    'summary-generation': {'mp3', 'mp4', 'wav', 'avi', 'mov'}
}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['PROCESSED_FOLDER'] = PROCESSED_FOLDER
app.config['IMAGE_FOLDER'] = IMAGE_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB max file size

# Create folders if they don't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PROCESSED_FOLDER, exist_ok=True)
os.makedirs(IMAGE_FOLDER, exist_ok=True)

def allowed_file(filename, service):
    """Check if file extension is allowed for the service"""
    if '.' not in filename:
        return False
    extension = filename.rsplit('.', 1)[1].lower()
    return extension in ALLOWED_EXTENSIONS.get(service, set())

def generate_unique_filename(original_filename):
    """Generate unique filename with timestamp and UUID"""
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    unique_id = str(uuid.uuid4())[:8]
    name, ext = os.path.splitext(original_filename)
    return f"{timestamp}{unique_id}{secure_filename(name)}{ext}"

# --- NOU: Funcția de parsare pentru fișiere PowerPoint ---
def parse_presentation(filepath):
    """
    Extrage text, imagini și metadate dintr-un fișier .pptx.
    """
    prs = Presentation(filepath)
    presentation_data = {
        "type": "ppt",
        "classification": "Prezentare (Clasificare Automată)",
        "presentation": os.path.basename(filepath),
        "totalSlides": len(prs.slides),
        "paragraphs": [],
        "images": []
    }

    current_subject = "Nespecificat"
    current_section = "General" # Secțiunea poate fi o logică mai complexă

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1

        # Heuristică simplă pentru a determina subiectul și secțiunea
        try:
            # Consideră titlul slide-ului ca fiind subiectul
            if slide.shapes.title:
                current_subject = slide.shapes.title.text.strip()
            # O logică mai avansată ar putea căuta slide-uri cu layout de "Section Header"
        except Exception:
            # Slide-ul poate să nu aibă un titlu
            pass

        for shape in slide.shapes:
            # Extrage textul din forme
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        para_data = {
                            "text": paragraph.text.strip(),
                            "slide": slide_number,
                            "section": current_section,
                            "subject": current_subject,
                            "level": f"bullet-{paragraph.level}",
                            "listType": "bullet" if paragraph.level > 0 else None,
                            "tableCell": None
                        }
                        presentation_data["paragraphs"].append(para_data)

            # Extrage textul din tabele
            if shape.has_table:
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            cell_data = {
                                "text": cell.text.strip(),
                                "slide": slide_number,
                                "section": current_section,
                                "subject": current_subject,
                                "level": "cell",
                                "listType": None,
                                "tableCell": {"row": r_idx + 1, "col": c_idx + 1}
                            }
                            presentation_data["paragraphs"].append(cell_data)

            # Extrage imaginile
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_filename = f"slide{slide_number}_{uuid.uuid4().hex[:8]}.{image.ext}"
                image_path = os.path.join(app.config['IMAGE_FOLDER'], image_filename)

                with open(image_path, 'wb') as img_file:
                    img_file.write(image_bytes)

                image_data = {
                    "url": f"/api/images/{image_filename}", # URL relativ pentru acces
                    "slide": slide_number,
                    "section": current_section,
                    "subject": current_subject,
                    "altText": shape.name or f"Imagine de pe slide-ul {slide_number}",
                    "position": {"left": shape.left.pt, "top": shape.top.pt},
                    "size": {"width": shape.width.pt, "height": shape.height.pt}
                }
                presentation_data["images"].append(image_data)

    return presentation_data

@app.route('/api/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        'status': 'online',
        'message': 'Backend server is running',
        'timestamp': datetime.now().isoformat()
    }), 200

@app.route('/api/document-parse', methods=['POST'])
def document_parse():
    """Process and classify documents"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        service = request.form.get('service', 'document-parse')
        
        if not allowed_file(file.filename, service):
            return jsonify({'error': 'File type not allowed'}), 400
        
        # Save uploaded file
        filename = generate_unique_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        # --- LOGICĂ ACTUALIZATĂ ---
        result = {}
        file_ext = filename.rsplit('.', 1)[1].lower()

        if file_ext in ['ppt', 'pptx']:
            # NOTĂ: Librăria suportă doar .pptx. Fișierele .ppt vechi vor da eroare.
            if file_ext == 'ppt':
                 return jsonify({'error': 'Formatul .ppt nu este suportat. Vă rugăm salvați fișierul ca .pptx.'}), 400
            result = parse_presentation(filepath)
        else:
            # Păstrăm simularea pentru alte tipuri de fișiere deocamdată
            result = {
                'type': 'pdf/doc/image',
                'document': file.filename,
                'status': 'success',
                'message': 'Procesarea reală pentru acest tip de fișier nu este încă implementată.'
            }

        return jsonify(result), 200
        
    except Exception as e:
        app.logger.error(f"Error in document_parse: {e}")
        return jsonify({'error': str(e)}), 500

# --- NOU: Ruta pentru a servi imaginile extrase ---
@app.route('/api/images/<path:filename>', methods=['GET'])
def serve_image(filename):
    """Serve extracted images"""
    try:
        return send_file(os.path.join(app.config['IMAGE_FOLDER'], filename))
    except FileNotFoundError:
        return jsonify({'error': 'Image not found'}), 404

# ... (restul rutelor API rămân neschimbate) ...
@app.route('/api/translation', methods=['POST'])
def translation():
    """Translate documents"""
    # ... (cod existent) ...
    return jsonify({'message': 'Endpoint not fully implemented'}), 200

@app.route('/api/audio-translation', methods=['POST'])
def audio_translation():
    """Translate audio files"""
    # ... (cod existent) ...
    return jsonify({'message': 'Endpoint not fully implemented'}), 200

@app.route('/api/video-translation', methods=['POST'])
def video_translation():
    """Translate video files"""
    # ... (cod existent) ...
    return jsonify({'message': 'Endpoint not fully implemented'}), 200

@app.route('/api/video-subtitle', methods=['POST'])
def video_subtitle():
    """Generate subtitles for videos"""
    # ... (cod existent) ...
    return jsonify({'message': 'Endpoint not fully implemented'}), 200

@app.route('/api/video-audio-replace', methods=['POST'])
def video_audio_replace():
    """Replace audio in video"""
    # ... (cod existent) ...
    return jsonify({'message': 'Endpoint not fully implemented'}), 200

@app.route('/api/summary-generation', methods=['POST'])
def summary_generation():
    """Generate summary for audio/video"""
    # ... (cod existent) ...
    return jsonify({'message': 'Endpoint not fully implemented'}), 200

@app.route('/api/download/<path:filename>', methods=['GET'])
def download_file(filename):
    """Download processed file"""
    try:
        # Check in both upload and processed folders
        for folder in [app.config['PROCESSED_FOLDER'], app.config['UPLOAD_FOLDER']]:
            filepath = os.path.join(folder, filename)
            if os.path.exists(filepath):
                return send_file(filepath, as_attachment=True)
        
        return jsonify({'error': 'File not found'}), 404
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File is too large. Maximum size is 500MB'}), 413

@app.errorhandler(500)
def internal_error(e):
    return jsonify({'error': 'Internal server error'}), 500

if _name_ == '_main_':
    print("=" * 60)
    print("🚀 Document Processing Backend Server")
    print("=" * 60)
    print(f"📂 Upload folder: {os.path.abspath(UPLOAD_FOLDER)}")
    print(f"📂 Processed folder: {os.path.abspath(PROCESSED_FOLDER)}")
    print(f"🖼 Image folder: {os.path.abspath(IMAGE_FOLDER)}")
    print(f"🌐 Server running on: http://localhost:5000")
    print(f"✅ Health check: http://localhost:5000/api/health")
    print("=" * 60)
    
    app.run(debug=True, host='0.0.0.0', port=5000)